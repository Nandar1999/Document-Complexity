import streamlit as st
import pdfplumber
import numpy as np
import pandas as pd
import fitz  # PyMuPDF (for images in PDF)
import docx
import pptx
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextBox

# Function to analyze table complexity
def calculate_table_complexity(table):
    num_rows = len(table)
    num_cols = max(len(row) for row in table) if table else 0
    total_cells = num_rows * num_cols
    empty_cells = sum(1 for row in table for cell in row if not cell or cell.strip() == "")
    merged_cell_ratio = empty_cells / total_cells if total_cells else 0
    row_lengths = [len(row) for row in table]
    row_length_variation = np.std(row_lengths) > 1
    nested_table = any(isinstance(cell, list) or len(str(cell).split()) > 15 for row in table for cell in row)
    header_count = sum(1 for row in table[:3] if all(cell and cell.isalpha() for cell in row if isinstance(cell, str)))
    word_counts = [len(str(cell).split()) for row in table for cell in row if cell]
    avg_word_count = np.mean(word_counts) if word_counts else 0
    high_density = avg_word_count > np.percentile(word_counts, 75) if word_counts else False
    
    score = (
        (0.4 if merged_cell_ratio > 0.2 else 0) +
        (0.3 if row_length_variation else 0) +
        (0.4 if nested_table else 0) +
        (0.3 if header_count > 1 else 0) +
        (0.2 if high_density else 0)
    )
    complexity_label = "Simple" if score <= 0.3 else "Moderate" if score <= 0.6 else "Complex"
    return score, complexity_label

# PDF Analysis Functions
def count_complex_tables_pdf(pdf_path):
    complex_table_count = 0
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            complex_table_count += sum(1 for table in tables if calculate_table_complexity(table)[1] == "Complex")
    return complex_table_count

def count_images_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    return sum(len(page.get_images(full=True)) for page in doc)

def analyze_layout_complexity_pdf(pdf_path):
    dense_paragraphs = 0
    num_columns = 1
    for page_layout in extract_pages(pdf_path):
        paragraph_count = sum(1 for element in page_layout if isinstance(element, LTTextBox) and len(element.get_text()) > 500)
        if paragraph_count > 5:
            num_columns = 2
        dense_paragraphs += paragraph_count
    return num_columns, dense_paragraphs

# DOCX Analysis Functions
def count_tables_docx(doc_path):
    doc = docx.Document(doc_path)
    return sum(1 for table in doc.tables if calculate_table_complexity([[cell.text for cell in row.cells] for row in table.rows])[1] == "Complex")

def count_images_docx(doc_path):
    return len(docx.Document(doc_path).inline_shapes)

def analyze_layout_complexity_docx(doc_path):
    doc = docx.Document(doc_path)
    dense_paragraphs = sum(1 for para in doc.paragraphs if len(para.text) > 500)
    num_columns = 1
    for section in doc.sections:
        if section._sectPr.xpath("./w:cols/w:col"):  # Check for multiple columns
            num_columns = max(num_columns, len(section._sectPr.xpath("./w:cols/w:col")))
    return num_columns, dense_paragraphs

# PPTX Analysis Functions
def count_tables_pptx(ppt_path):
    ppt = pptx.Presentation(ppt_path)
    return sum(1 for slide in ppt.slides for shape in slide.shapes if shape.has_table and calculate_table_complexity([[cell.text for cell in row.cells] for row in shape.table.rows])[1] == "Complex")

def count_images_pptx(ppt_path):
    return sum(1 for slide in pptx.Presentation(ppt_path).slides for shape in slide.shapes if shape.shape_type == 13)

def analyze_layout_complexity_pptx(ppt_path):
    ppt = pptx.Presentation(ppt_path)
    dense_paragraphs = sum(1 for slide in ppt.slides for shape in slide.shapes if shape.has_text_frame and len(shape.text_frame.text) > 500)
    return (1 if dense_paragraphs < 5 else 2), dense_paragraphs

# Main Complexity Calculation
def calculate_document_complexity(file_path, file_type):
    if file_type == "pdf":
        num_complex_tables = count_complex_tables_pdf(file_path)
        num_images = count_images_pdf(file_path)
        num_columns, dense_paragraphs = analyze_layout_complexity_pdf(file_path)
    elif file_type == "docx":
        num_complex_tables = count_tables_docx(file_path)
        num_images = count_images_docx(file_path)
        num_columns, dense_paragraphs = analyze_layout_complexity_docx(file_path)
    elif file_type == "pptx":
        num_complex_tables = count_tables_pptx(file_path)
        num_images = count_images_pptx(file_path)
        num_columns, dense_paragraphs = analyze_layout_complexity_pptx(file_path)
        dense_paragraphs = 0
    
    table_score = num_complex_tables * 20
    layout_score = (num_columns * 20 + dense_paragraphs * 5) if (num_columns > 1 or dense_paragraphs > 0) else 0
    image_score = num_images * 10
    final_score = table_score + layout_score + image_score
    complexity_level = "High" if final_score > 100 else "Medium" if final_score > 50 else "Low"
    
    return {
        "Complex Tables Found": num_complex_tables,
        "Columns Detected": num_columns,
        "Dense Paragraphs": dense_paragraphs,
        "Images Found": num_images,
        "Final Complexity Score": final_score,
        "Complexity Level": complexity_level
    }

# Streamlit UI
st.set_page_config(page_title="Document Complexity Detector", layout="wide")
st.title("📄 Document Complexity Detector")
st.write("Upload a PDF, DOCX, or PPTX to detect its complexity based on **complex tables, layout structure, and images**.")

doc_type = st.selectbox("Select Document Type", ["PDF", "DOCX", "PPTX"])

uploaded_file = st.file_uploader("📂 Upload your document", type=[doc_type.lower()])

if uploaded_file:
    file_path = f"temp_uploaded.{doc_type.lower()}"
    with open(file_path, "wb") as f:
        f.write(uploaded_file.read())
    st.success(f"✅ {doc_type.upper()} file uploaded successfully!")
    with st.spinner("🔍 Analyzing document..."):
        complexity_results = calculate_document_complexity(file_path, doc_type.lower())
    st.subheader("📊 Complexity Analysis Report")
    for key, value in complexity_results.items():
        st.metric(label=key, value=value)
    st.success("✅ Analysis Complete!")
